"""
create_presentation.py - Generate project presentation PowerPoint
=================================================================
Run:  python create_presentation.py
Output: outputs/BTC_Sentiment_Prediction_Presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
DARK_BG = RGBColor(0x1B, 0x1B, 0x2F)
ACCENT  = RGBColor(0xF7, 0x93, 0x1E)  # orange
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT   = RGBColor(0xCC, 0xCC, 0xCC)
GREEN   = RGBColor(0x4E, 0xC9, 0xB0)
RED     = RGBColor(0xE0, 0x5E, 0x5E)
BLUE    = RGBColor(0x56, 0x9C, 0xD6)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_slide_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_paragraph(tf, text, font_size=16, color=WHITE, bold=False,
                  space_before=Pt(6), bullet=False, font_name="Calibri"):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    if space_before:
        p.space_before = space_before
    if bullet:
        p.level = 1
    return p


def add_image_safe(slide, path, left, top, width=None, height=None):
    if Path(path).exists():
        slide.shapes.add_picture(path, left, top, width, height)
        return True
    return False


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------
def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide)

    add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
                 "BTC/USD Price Prediction with News Sentiment",
                 font_size=36, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1), Inches(3.0), Inches(11), Inches(0.8),
                 "Does FinBERT news sentiment improve next-day return prediction\n"
                 "over a technical-only baseline?",
                 font_size=22, color=WHITE, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1), Inches(5.0), Inches(11), Inches(0.6),
                 "Machine Learning Project  |  April 2026",
                 font_size=18, color=LIGHT, alignment=PP_ALIGN.CENTER)

    # Placeholder for student info
    add_text_box(slide, Inches(1), Inches(5.8), Inches(11), Inches(0.8),
                 "[Student Name(s) & ID(s) - UPDATE THIS]",
                 font_size=16, color=ACCENT, alignment=PP_ALIGN.CENTER)


def slide_overview(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
                 "Project Overview", font_size=30, color=ACCENT, bold=True)

    tf = add_text_box(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(5.5),
                      "Research Question", font_size=20, color=GREEN, bold=True)
    add_paragraph(tf, 'Can lagged FinBERT news sentiment improve\n'
                      'next-day BTC/USD log-return prediction\n'
                      'vs. a technical-indicator-only model?', font_size=16, color=WHITE)

    add_paragraph(tf, "", font_size=10)
    add_paragraph(tf, "Approach: Two-Stage Model", font_size=20, color=GREEN, bold=True)
    add_paragraph(tf, "Stage 1 (Model A): LightGBM trained on 5 years\n"
                      "of 16 technical indicators (RSI, MACD, Bollinger\n"
                      "Bands, moving averages, volatility, volume, etc.)", font_size=15, color=WHITE)
    add_paragraph(tf, "Stage 2 (Model B): ElasticNet trained on the last\n"
                      "65 days to correct Model A's residuals using 11\n"
                      "FinBERT sentiment features from crypto news", font_size=15, color=WHITE)
    add_paragraph(tf, "Final Prediction = Model A + Model B correction", font_size=15, color=ACCENT, bold=True)

    # Right side - pipeline diagram as text
    tf2 = add_text_box(slide, Inches(6.8), Inches(1.3), Inches(5.8), Inches(5.5),
                       "Pipeline Architecture", font_size=20, color=GREEN, bold=True)
    steps = [
        ("Step 1:", "Load OHLCV data + collect news articles"),
        ("Step 2:", "Engineer 16 technical + 11 sentiment features"),
        ("Step 3:", "Score headlines with FinBERT transformer"),
        ("Step 4:", "Train LightGBM (Model A) on technicals"),
        ("Step 5:", "Extract residuals, train ElasticNet (Model B)"),
        ("Step 6:", "Walk-forward evaluation with expanding window"),
        ("Step 7:", "Backtest with stop-loss / take-profit rules"),
    ]
    for label, desc in steps:
        add_paragraph(tf2, f"{label}  {desc}", font_size=14, color=WHITE)


def slide_data(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
                 "Datasets & Features", font_size=30, color=ACCENT, bold=True)

    # Dataset A
    tf = add_text_box(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(2.5),
                      "Dataset A: BTC/USD OHLCV (Investing.com)", font_size=18, color=GREEN, bold=True)
    add_paragraph(tf, "1,827 daily rows: Apr 2021 - Apr 2026", font_size=15, color=WHITE)
    add_paragraph(tf, "Columns: Date, Price, Open, High, Low, Volume", font_size=15, color=WHITE)
    add_paragraph(tf, "Source: investing.com/crypto/bitcoin/historical-data", font_size=13, color=LIGHT)

    # Dataset B
    tf2 = add_text_box(slide, Inches(0.8), Inches(3.5), Inches(5.5), Inches(2.0),
                       "Dataset B: Crypto News (RSS + Pre-scraped)", font_size=18, color=GREEN, bold=True)
    add_paragraph(tf2, "75 days of news coverage (Jan - Apr 2026)", font_size=15, color=WHITE)
    add_paragraph(tf2, "Sources: CoinDesk, CoinTelegraph, Decrypt,\n"
                       "Bitcoin Magazine, CryptoSlate, The Block", font_size=15, color=WHITE)
    add_paragraph(tf2, "Scored with ProsusAI/FinBERT transformer", font_size=15, color=WHITE)

    # External datasets
    tf3 = add_text_box(slide, Inches(0.8), Inches(5.3), Inches(5.5), Inches(1.5),
                       "External Signals (Optional)", font_size=18, color=GREEN, bold=True)
    add_paragraph(tf3, "Google Trends: bitcoin, btc, bitcoin price", font_size=15, color=WHITE)
    add_paragraph(tf3, "Fear & Greed Index (Alternative.me)", font_size=15, color=WHITE)

    # Feature table on right
    tf4 = add_text_box(slide, Inches(6.8), Inches(1.3), Inches(5.8), Inches(5.5),
                       "Feature Groups", font_size=20, color=GREEN, bold=True)

    add_paragraph(tf4, "", font_size=6)
    add_paragraph(tf4, "Technical (16 features - Model A):", font_size=16, color=ACCENT, bold=True)
    add_paragraph(tf4, "log_return, return_lag1, hl_spread, oc_spread\n"
                       "vol_log, volume_change_pct, obv_change_pct\n"
                       "price_vs_sma20/50/200, rsi_14, macd_hist\n"
                       "bband_pos, atr_pct, vol_10d, vol_30d", font_size=13, color=WHITE)

    add_paragraph(tf4, "", font_size=6)
    add_paragraph(tf4, "Sentiment (11 features - Model B):", font_size=16, color=ACCENT, bold=True)
    add_paragraph(tf4, "sentiment_mean/std/max/min/count\n"
                       "positive_ratio, negative_ratio, sentiment_range\n"
                       "sentiment_mom_3d/7d, bull_bear_ratio", font_size=13, color=WHITE)

    add_paragraph(tf4, "", font_size=6)
    add_paragraph(tf4, "Leakage Control:", font_size=16, color=RED, bold=True)
    add_paragraph(tf4, "All external features lagged by 1 day\n"
                       "Last 65 days held out from Model A training", font_size=13, color=WHITE)


def slide_model_training(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
                 "Model Training", font_size=30, color=ACCENT, bold=True)

    # Model A - both models
    tf = add_text_box(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(5.5),
                      "Model A: LightGBM vs XGBoost (Regression)", font_size=18, color=GREEN, bold=True)
    add_paragraph(tf, "Target: next-day log return", font_size=15, color=WHITE)
    add_paragraph(tf, "Training: 1,762 rows (excl. 65-day news window)", font_size=15, color=WHITE)
    add_paragraph(tf, "Validation: 5-fold TimeSeriesSplit (no shuffle)", font_size=15, color=WHITE)
    add_paragraph(tf, "Early stopping: patience=50 iterations", font_size=15, color=WHITE)

    add_paragraph(tf, "", font_size=6)
    add_paragraph(tf, "Cross-Validation Comparison:", font_size=16, color=ACCENT, bold=True)
    add_paragraph(tf, "LightGBM  CV RMSE: 0.02630 (iter=1-9)", font_size=14, color=WHITE)
    add_paragraph(tf, "XGBoost   CV RMSE: 0.02626 (iter=1-31)", font_size=14, color=GREEN)
    add_paragraph(tf, "Winner: XGBoost (lower RMSE by 0.000035)", font_size=14, color=GREEN, bold=True)

    add_paragraph(tf, "", font_size=6)
    add_paragraph(tf, "Backtest Results (Tech Only):", font_size=16, color=ACCENT, bold=True)
    add_paragraph(tf, "LightGBM: -4.62% return, Sharpe -0.52", font_size=14, color=RED)
    add_paragraph(tf, "XGBoost:  +6.00% return, Sharpe +1.07", font_size=14, color=GREEN)
    add_paragraph(tf, "Buy&Hold: +6.97% return, Sharpe +1.07", font_size=14, color=LIGHT)
    add_paragraph(tf, "", font_size=4)
    add_paragraph(tf, "XGBoost nearly matches Buy & Hold while\n"
                      "LightGBM loses money on the same data.", font_size=13, color=LIGHT)

    # Model B
    tf2 = add_text_box(slide, Inches(6.8), Inches(1.3), Inches(5.8), Inches(5.5),
                       "Model B: ElasticNet (Residual Corrector)", font_size=20, color=GREEN, bold=True)
    add_paragraph(tf2, "Target: Model A residuals (what A got wrong)", font_size=15, color=WHITE)
    add_paragraph(tf2, "Training: 64 rows with actual news coverage", font_size=15, color=WHITE)
    add_paragraph(tf2, "CV: TimeSeriesSplit (2-3 folds, auto-tuned)", font_size=15, color=WHITE)

    add_paragraph(tf2, "", font_size=8)
    add_paragraph(tf2, "Why ElasticNet?", font_size=16, color=ACCENT, bold=True)
    add_paragraph(tf2, "Only 64 samples = needs strong regularization\n"
                       "L1 (Lasso): automatic feature selection\n"
                       "L2 (Ridge): handles correlated sentiment features\n"
                       "ElasticNetCV auto-tunes alpha & l1_ratio", font_size=14, color=WHITE)

    add_paragraph(tf2, "", font_size=8)
    add_paragraph(tf2, "Result:", font_size=16, color=ACCENT, bold=True)
    add_paragraph(tf2, "Best alpha=0.032, l1_ratio=0.90", font_size=15, color=WHITE)
    add_paragraph(tf2, "Non-zero sentiment features: 0 / 11", font_size=15, color=RED)
    add_paragraph(tf2, "Sentiment was fully regularized to zero\n"
                       "(too few news samples for signal to emerge)", font_size=14, color=LIGHT)


def slide_evaluation_chart(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
                 "Evaluation: Predicted vs Actual Returns", font_size=30, color=ACCENT, bold=True)

    img = "outputs/btc_two_stage_evaluation_twostage_lightgbm.png"
    if not add_image_safe(slide, img, Inches(0.5), Inches(1.1), width=Inches(12.3)):
        add_text_box(slide, Inches(2), Inches(3), Inches(8), Inches(1),
                     f"[Chart not found: {img}]", font_size=20, color=RED, alignment=PP_ALIGN.CENTER)


def slide_results_table(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
                 "Model Comparison Results", font_size=30, color=ACCENT, bold=True)

    # Walk-forward metrics
    tf = add_text_box(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(1.0),
                      "Walk-Forward Evaluation (last 50 days after 15-day warm-up)",
                      font_size=20, color=GREEN, bold=True)

    add_paragraph(tf, "", font_size=6)

    # Table as formatted text
    header = f"{'Metric':<30} {'Model A (Tech Only)':<22} {'Ensemble (A+B)':<22}"
    add_paragraph(tf, header, font_size=15, color=ACCENT, bold=True, font_name="Courier New")

    rows = [
        ("RMSE",              "0.024682",  "0.025206"),
        ("MAE",               "0.019022",  "0.019406"),
        ("Direction Accuracy", "46.00%",    "50.00%"),
        ("Dir.Acc Lift",      "--",        "+4.0 pp"),
    ]
    for label, a, b in rows:
        line = f"{label:<30} {a:<22} {b:<22}"
        add_paragraph(tf, line, font_size=14, color=WHITE, font_name="Courier New")

    # Ablation table
    tf2 = add_text_box(slide, Inches(0.8), Inches(4.2), Inches(11.5), Inches(3.0),
                       "Ablation Study: Feature Set Comparison (all LightGBM)",
                       font_size=20, color=GREEN, bold=True)

    header2 = f"{'Variant':<28} {'Features':<10} {'RMSE':<10} {'Dir.Acc':<10} {'Backtest Return':<16}"
    add_paragraph(tf2, header2, font_size=14, color=ACCENT, bold=True, font_name="Courier New")

    ablation_rows = [
        ("A0: Tech Only",          "16", "0.0334", "43.1%", "-4.62%"),
        ("A1: Tech + Sentiment",   "27", "0.0334", "43.1%", "-4.62%"),
        ("A2: Tech + Trends",      "20", "0.0338", "44.6%", "-7.08%"),
        ("A3: Tech + FGI",         "18", "0.0335", "41.5%", "-7.78%"),
        ("A4: Tech + All",         "33", "0.0343", "43.1%", "-7.14%"),
        ("Two-Stage (A+B)",        "27", "0.0252", "50.0%", "+0.88%"),
    ]
    for label, feats, rmse, dacc, ret in ablation_rows:
        c = GREEN if "Two-Stage" in label else WHITE
        line = f"{label:<28} {feats:<10} {rmse:<10} {dacc:<10} {ret:<16}"
        add_paragraph(tf2, line, font_size=13, color=c, font_name="Courier New")

    add_paragraph(tf2, "", font_size=6)
    add_paragraph(tf2, "A1 vs A0 sentiment lift: 0.00 pp (sentiment adds nothing in single-stage LightGBM)",
                  font_size=14, color=LIGHT)
    add_paragraph(tf2, "Two-Stage ensemble is the only profitable strategy (+0.88% vs Buy&Hold +6.97%)",
                  font_size=14, color=LIGHT)


def slide_backtest(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
                 "Backtest: Equity Curves", font_size=30, color=ACCENT, bold=True)

    img = "outputs/btc_backtest_equity_curve_twostage_lightgbm.png"
    if not add_image_safe(slide, img, Inches(0.3), Inches(1.0), width=Inches(7.0)):
        add_text_box(slide, Inches(1), Inches(2.5), Inches(5), Inches(1),
                     f"[Chart not found: {img}]", font_size=18, color=RED)

    # Summary table on right
    tf = add_text_box(slide, Inches(7.6), Inches(1.0), Inches(5.2), Inches(5.5),
                      "Backtest Summary (50 days)", font_size=20, color=GREEN, bold=True)

    add_paragraph(tf, "Stop-Loss: -2%  |  Take-Profit: +4%", font_size=14, color=LIGHT)
    add_paragraph(tf, "Trading Fee: 0.1% per trade", font_size=14, color=LIGHT)

    add_paragraph(tf, "", font_size=8)
    add_paragraph(tf, "Model A (Tech Only)", font_size=16, color=BLUE, bold=True)
    add_paragraph(tf, "Return: -4.62%  |  49 trades  |  Win: 42.9%\n"
                      "Sharpe: -0.52  |  Max DD: -13.3%", font_size=14, color=WHITE)

    add_paragraph(tf, "", font_size=8)
    add_paragraph(tf, "Ensemble (A+B)", font_size=16, color=ACCENT, bold=True)
    add_paragraph(tf, "Return: +0.88%  |  3 trades  |  Win: 33.3%\n"
                      "Sharpe: +0.51  |  Max DD: -2.8%", font_size=14, color=WHITE)

    add_paragraph(tf, "", font_size=8)
    add_paragraph(tf, "Buy & Hold BTC", font_size=16, color=LIGHT, bold=True)
    add_paragraph(tf, "Return: +6.97%  |  Sharpe: +1.07\n"
                      "Max DD: -11.9%", font_size=14, color=WHITE)

    add_paragraph(tf, "", font_size=8)
    add_paragraph(tf, "Key Insight: Ensemble trades rarely (3 vs 49)\n"
                      "but avoids large losses. Conservative but\n"
                      "risk-managed.", font_size=14, color=GREEN)


def slide_direction_classification(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
                 "Direction Classification (Up/Down)", font_size=30, color=ACCENT, bold=True)

    tf = add_text_box(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(1.0),
                      "Binary classification: will BTC go up or down tomorrow?",
                      font_size=18, color=WHITE)

    add_paragraph(tf, "", font_size=8)

    header = f"{'Variant':<28} {'Model':<22} {'Accuracy':<10} {'Precision':<10} {'Recall':<10} {'F1':<10}"
    add_paragraph(tf, header, font_size=14, color=ACCENT, bold=True, font_name="Courier New")

    cls_rows = [
        ("C0: Tech Only",         "LightGBM",           "52.3%", "43.5%", "35.7%", "39.2%"),
        ("C0: Tech Only",         "Logistic Regression", "55.4%", "47.4%", "32.1%", "38.3%"),
        ("C1: Tech + Sentiment",  "LightGBM",           "50.8%", "42.3%", "39.3%", "40.7%"),
        ("C2: Tech + All",        "LightGBM",           "50.8%", "38.9%", "25.0%", "30.4%"),
    ]
    for variant, model, acc, prec, rec, f1 in cls_rows:
        c = GREEN if "40.7" in f1 else WHITE
        line = f"{variant:<28} {model:<22} {acc:<10} {prec:<10} {rec:<10} {f1:<10}"
        add_paragraph(tf, line, font_size=13, color=c, font_name="Courier New")

    add_paragraph(tf, "", font_size=10)
    add_paragraph(tf, "Best F1 = 40.7% with Tech + Sentiment (LightGBM classifier)",
                  font_size=16, color=GREEN, bold=True)
    add_paragraph(tf, "Sentiment marginally improves recall (+3.6pp) in classification,\n"
                      "even though it had no effect in regression (A1 vs A0).",
                  font_size=15, color=WHITE)


def slide_improved_overview(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
                 "Improved Model: Direction Classifier + Confidence Threshold",
                 font_size=28, color=ACCENT, bold=True)

    tf = add_text_box(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(5.5),
                      "3 Key Improvements", font_size=20, color=GREEN, bold=True)

    add_paragraph(tf, "", font_size=6)
    add_paragraph(tf, "1. Direction Classification (not regression)", font_size=17, color=ACCENT, bold=True)
    add_paragraph(tf, "Predict UP or DOWN instead of exact return.\n"
                      "Trading decisions only need direction, not\n"
                      "magnitude. Simpler, more learnable task.", font_size=14, color=WHITE)

    add_paragraph(tf, "", font_size=8)
    add_paragraph(tf, "2. Single-Stage Model (tech + sentiment)", font_size=17, color=ACCENT, bold=True)
    add_paragraph(tf, "Train ONE model on ALL 1,762 rows with both\n"
                      "feature types. LightGBM handles sparse\n"
                      "sentiment naturally (tree-based splits).", font_size=14, color=WHITE)

    add_paragraph(tf, "", font_size=8)
    add_paragraph(tf, "3. Confidence Threshold Trading", font_size=17, color=ACCENT, bold=True)
    add_paragraph(tf, "Only trade when P(up) >= threshold.\n"
                      "Filters low-conviction signals, dramatically\n"
                      "reduces bad trades.", font_size=14, color=WHITE)

    # Right side - threshold results
    tf2 = add_text_box(slide, Inches(6.8), Inches(1.3), Inches(5.8), Inches(5.5),
                       "Threshold Analysis (Tech + Sentiment)", font_size=20, color=GREEN, bold=True)

    header = f"{'Threshold':>10}  {'Trades':>7}  {'Accuracy':>9}  {'F1':>7}"
    add_paragraph(tf2, header, font_size=15, color=ACCENT, bold=True, font_name="Courier New")

    rows = [
        ("50%", "65", "60.0%", "0.552"),
        ("52%", "36", "52.8%", "0.485"),
        ("55%", "10", "70.0%", "0.769"),
        ("58%", "4", "75.0%", "0.857"),
    ]
    for t, trades, acc, f1 in rows:
        c = GREEN if "70" in acc or "75" in acc else WHITE
        line = f"{t:>10}  {trades:>7}  {acc:>9}  {f1:>7}"
        add_paragraph(tf2, line, font_size=14, color=c, font_name="Courier New")

    add_paragraph(tf2, "", font_size=10)
    add_paragraph(tf2, "Key insight: at 55% threshold, accuracy jumps\n"
                       "from 60% to 70% while trading only 10/65 days.\n"
                       "Quality over quantity.", font_size=15, color=GREEN)

    add_paragraph(tf2, "", font_size=8)
    add_paragraph(tf2, "At 58% threshold: 75% accuracy, F1=0.857\n"
                       "(only 4 trades but very high conviction)", font_size=14, color=LIGHT)


def slide_improved_backtest(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
                 "Improved Backtest Results", font_size=30, color=ACCENT, bold=True)

    img = "outputs/improved_equity_comparison.png"
    if not add_image_safe(slide, img, Inches(0.3), Inches(1.0), width=Inches(7.5)):
        add_text_box(slide, Inches(1), Inches(2.5), Inches(5), Inches(1),
                     f"[Chart not found: {img}]", font_size=18, color=RED)

    tf = add_text_box(slide, Inches(8.0), Inches(1.0), Inches(5.0), Inches(6.0),
                      "Strategy Comparison", font_size=18, color=GREEN, bold=True)

    add_paragraph(tf, "", font_size=6)
    add_paragraph(tf, "Old: Tech Regression", font_size=15, color=BLUE, bold=True)
    add_paragraph(tf, "Return: -20.4%  |  64 trades\nSharpe: -2.35  |  Max DD: -25.7%", font_size=13, color=WHITE)

    add_paragraph(tf, "", font_size=6)
    add_paragraph(tf, "New: Classifier (50% threshold)", font_size=15, color=RED, bold=True)
    add_paragraph(tf, "Return: -2.6%  |  29 trades\nSharpe: -0.30  |  Max DD: -9.3%", font_size=13, color=WHITE)

    add_paragraph(tf, "", font_size=6)
    add_paragraph(tf, "New: Classifier (55% threshold)", font_size=15, color=ACCENT, bold=True)
    add_paragraph(tf, "Return: +6.51%  |  8 trades\nSharpe: +1.67  |  Max DD: -4.4%", font_size=13, color=GREEN)

    add_paragraph(tf, "", font_size=6)
    add_paragraph(tf, "Buy & Hold BTC", font_size=15, color=LIGHT, bold=True)
    add_paragraph(tf, "Return: -6.1%  |  Sharpe: -0.20\nMax DD: -17.0%", font_size=13, color=WHITE)

    add_paragraph(tf, "", font_size=10)
    add_paragraph(tf, "55% threshold BEATS Buy & Hold:", font_size=16, color=GREEN, bold=True)
    add_paragraph(tf, "+6.51% vs -6.10%\nSharpe 1.67 vs -0.20\nMax DD -4.4% vs -17.0%", font_size=14, color=GREEN)

    add_paragraph(tf, "", font_size=6)
    add_paragraph(tf, "Improvement: +26.95 pp over old model", font_size=14, color=ACCENT, bold=True)


def slide_feature_importance(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
                 "Feature Importance: Direction Classifier", font_size=30, color=ACCENT, bold=True)

    img = "outputs/improved_feature_importance.png"
    if not add_image_safe(slide, img, Inches(0.3), Inches(1.0), width=Inches(7.5)):
        add_text_box(slide, Inches(1), Inches(2.5), Inches(5), Inches(1),
                     f"[Chart not found: {img}]", font_size=18, color=RED)

    tf = add_text_box(slide, Inches(8.0), Inches(1.0), Inches(5.0), Inches(5.5),
                      "Top Drivers", font_size=20, color=GREEN, bold=True)

    add_paragraph(tf, "", font_size=6)
    add_paragraph(tf, "Most important features:", font_size=16, color=ACCENT, bold=True)
    add_paragraph(tf, "1. return_lag1 (37) - yesterday's return\n"
                      "2. price_vs_sma20 (30) - trend position\n"
                      "3. obv_change_pct (24) - volume flow\n"
                      "4. vol_10d (18) - short-term volatility\n"
                      "5. volume_change_pct (18) - volume shift",
                  font_size=14, color=WHITE)

    add_paragraph(tf, "", font_size=10)
    add_paragraph(tf, "Sentiment features:", font_size=16, color=RED, bold=True)
    add_paragraph(tf, "sentiment_mom_7d, sentiment_mean,\n"
                      "sentiment_std, sentiment_max\n"
                      "all show near-zero importance.",
                  font_size=14, color=WHITE)

    add_paragraph(tf, "", font_size=8)
    add_paragraph(tf, "Why? Only 72/1,827 days have real news.\n"
                      "96% of rows have zero/forward-filled\n"
                      "sentiment = no signal for tree splits.",
                  font_size=14, color=LIGHT)

    add_paragraph(tf, "", font_size=8)
    add_paragraph(tf, "Conclusion: The improvement comes from\n"
                      "classification + confidence threshold,\n"
                      "not from sentiment features.",
                  font_size=15, color=GREEN, bold=True)


def slide_limitations(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
                 "What We Learned / Limitations", font_size=30, color=ACCENT, bold=True)

    tf = add_text_box(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(5.5),
                      "Key Findings", font_size=22, color=GREEN, bold=True)

    findings = [
        "Direction classification + confidence threshold\n"
        "is the key breakthrough: +6.51% return, Sharpe\n"
        "1.67, beating Buy & Hold (-6.1%)",

        "Confidence threshold is critical: at 55%,\n"
        "accuracy jumps from 60% to 70% by filtering\n"
        "out low-conviction predictions (8 vs 64 trades)",

        "Sentiment adds NO value in any model config\n"
        "(A1 = A0 exactly, feature importance near zero)\n"
        "due to sparse news coverage (72/1,827 days)",

        "Max drawdown dramatically improved:\n"
        "-4.4% (new) vs -25.7% (old regression)\n"
        "vs -17.0% (Buy & Hold)",

        "Improvement of +26.95 pp over old regression\n"
        "model (-20.4% to +6.51%) by changing the\n"
        "problem framing, not the features",
    ]
    for f in findings:
        add_paragraph(tf, f, font_size=14, color=WHITE, space_before=Pt(12))

    tf2 = add_text_box(slide, Inches(6.8), Inches(1.3), Inches(5.8), Inches(5.5),
                       "Limitations", font_size=22, color=RED, bold=True)

    limitations = [
        "News coverage too sparse: 72 out of 1,827 days\n"
        "means sentiment features have near-zero impact",

        "FinBERT was pretrained on financial text, but\n"
        "crypto-specific slang may not score well",

        "Short backtest window (65 days) makes\n"
        "performance metrics statistically unreliable",

        "No transaction cost modeling beyond flat 0.1% fee;\n"
        "slippage, spread, and market impact are ignored",

        "High confidence threshold (55%+) trades very\n"
        "rarely (8 out of 65 days) - small sample size\n"
        "for statistical significance",

        "Model may be overfitting to recent regime;\n"
        "needs longer out-of-sample testing",
    ]
    for l in limitations:
        add_paragraph(tf2, l, font_size=14, color=WHITE, space_before=Pt(12))


def slide_summary(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
                 "Summary & Conclusions", font_size=30, color=ACCENT, bold=True)

    tf = add_text_box(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(2.5),
                      "Answer to Research Question", font_size=22, color=GREEN, bold=True)

    add_paragraph(tf, "Sentiment does NOT improve prediction (A1 = A0, feature importance near zero).",
                  font_size=18, color=RED, bold=True)
    add_paragraph(tf, "However, switching from regression to direction classification with a 55%\n"
                      "confidence threshold produces +6.51% return (Sharpe 1.67), beating Buy & Hold (-6.1%).\n"
                      "Improvement of +26.95 pp over the original regression model.",
                  font_size=17, color=GREEN)
    add_paragraph(tf, "The breakthrough is problem framing (classification + confidence filtering),\n"
                      "not the features. Sentiment needs denser news coverage (72/1,827 days is too sparse).",
                  font_size=16, color=LIGHT)

    # Predictions / Future Work
    tf2 = add_text_box(slide, Inches(0.8), Inches(4.5), Inches(5.5), Inches(2.5),
                       "Predictions / Future Work", font_size=20, color=GREEN, bold=True)
    future = [
        "More news sources (Twitter/Reddit) could\nprovide daily coverage for all trading days",
        "Multi-day horizons (3-5 day returns) would\nreduce noise in the target variable",
        "Adaptive threshold: adjust confidence cutoff\nbased on market volatility regime",
        "Longer out-of-sample test (6-12 months)\nto validate the 55% threshold strategy",
    ]
    for f in future:
        add_paragraph(tf2, f, font_size=14, color=WHITE, space_before=Pt(10))

    # Data sources
    tf3 = add_text_box(slide, Inches(6.8), Inches(4.5), Inches(5.8), Inches(2.5),
                       "Data Sources", font_size=20, color=GREEN, bold=True)
    sources = [
        "BTC/USD OHLCV: investing.com\n  /crypto/bitcoin/historical-data",
        "News: CoinDesk, CoinTelegraph, Decrypt,\n  Bitcoin Magazine, CryptoSlate, The Block (RSS)",
        "Sentiment Model: ProsusAI/FinBERT\n  (HuggingFace transformers)",
        "Fear & Greed: alternative.me API",
        "Google Trends: pytrends library",
    ]
    for s in sources:
        add_paragraph(tf3, s, font_size=13, color=WHITE, space_before=Pt(8))


def slide_thank_you(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(1), Inches(2.5), Inches(11), Inches(1),
                 "Thank You", font_size=44, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.8),
                 "Questions?", font_size=28, color=WHITE, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.8),
                 "Code & data available on request\n"
                 "[GitHub/OneDrive link - UPDATE THIS]",
                 font_size=16, color=LIGHT, alignment=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------
def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs)                    # 1.  Title
    slide_overview(prs)                 # 2.  Overview + Architecture
    slide_data(prs)                     # 3.  Datasets & Features
    slide_model_training(prs)           # 4.  Model Training Details
    slide_evaluation_chart(prs)         # 5.  Evaluation Chart (original)
    slide_results_table(prs)            # 6.  Results Comparison Table
    slide_backtest(prs)                 # 7.  Backtest Equity Curves (original)
    slide_improved_overview(prs)        # 8.  Improved Model: 3 key changes
    slide_improved_backtest(prs)        # 9.  Improved Backtest Results
    slide_feature_importance(prs)       # 10. Feature Importance Analysis
    slide_limitations(prs)              # 11. Limitations & Learnings
    slide_summary(prs)                  # 12. Summary & Conclusions
    slide_thank_you(prs)                # 13. Thank You

    out_path = Path("outputs/BTC_Sentiment_Prediction_Presentation.pptx")
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    print(f"Presentation saved -> {out_path}")
    print(f"  {len(prs.slides)} slides")
    print("\n  IMPORTANT: Update slide 1 with your name(s) & student ID(s)")
    print("  IMPORTANT: Update last slide with your GitHub/OneDrive link")


if __name__ == "__main__":
    main()
